"""PowerPoint parser using python-pptx (offline)."""
from pathlib import Path

from ..models import DocumentFormat, ParsedDocument
from .base import BaseParser


class PptxParser(BaseParser):
    supported_formats = [DocumentFormat.PPTX]

    def parse(self, path: Path) -> ParsedDocument:
        try:
            from pptx import Presentation

            prs = Presentation(str(path))
            lines: list[str] = []
            for i, slide in enumerate(prs.slides, 1):
                lines.append(f"--- Slide {i} ---")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            txt = para.text.strip()
                            if txt:
                                lines.append(txt)
            raw_text = "\n".join(lines)
            page_count = len(prs.slides)
        except ImportError:
            raw_text = f"[python-pptx not installed — cannot parse {path.name}]"
            page_count = 0
        except Exception as exc:  # noqa: BLE001
            raw_text = f"[PPTX parse error: {exc}]"
            page_count = 0

        return ParsedDocument(
            source_path=str(path),
            format=DocumentFormat.PPTX,
            title=path.stem,
            raw_text=raw_text,
            page_count=page_count,
        )
